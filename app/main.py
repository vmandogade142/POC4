from flask import Flask, request, render_template, jsonify, redirect, url_for, session
from models.vector_store import VectorStore
from services.storage_service import S3Storage
from services.llm_service import LLMService
from config import Config
import os
from langchain.document_loaders import TextLoader, PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document
import tempfile
import logging
import pandas as pd
from pptx import Presentation
from docx import Document as DocxDocument

# ----------------------------------------------------------------------------------------
# Initializing Flask App
app = Flask(__name__)
app.secret_key = Config.SECRET_KEY  # Use the secret key from the config file

vector_store = VectorStore(Config.VECTOR_DB_PATH)
storage_service = S3Storage()
llm_service = LLMService(vector_store)

# ----------------------------------------------------------------------------------------
# Configure logging
logging.basicConfig(level=logging.DEBUG)
logger = logging.getLogger(__name__)

# ----------------------------------------------------------------------------------------
# Dummy user database (for demonstration)
users = {
    "admin": {"password": "admin123", "role": "admin"},
    "user": {"password": "user123", "role": "user"}
}

# ----------------------------------------------------------------------------------------
# Home Route
@app.route('/')
def home():
    return render_template('index.html')  # Ensure index.html exists in the templates folder

# ----------------------------------------------------------------------------------------
# Authentication Routes
@app.route('/login', methods=['POST'])
def login():
    data = request.json
    username = data.get('username')
    password = data.get('password')
    user = users.get(username)

    if user and user['password'] == password:
        session['user'] = username
        session['role'] = user['role']
        return jsonify({'success': True, 'role': user['role']})
    return jsonify({'success': False, 'error': 'Invalid Credentials'}), 401

@app.route('/logout', methods=['POST'])
def logout():
    session.clear()
    return jsonify({'success': True})

# ----------------------------------------------------------------------------------------
# Protected Route
@app.route('/dashboard')
def dashboard():
    if 'user' not in session:
        return redirect(url_for('home'))  # Redirect to home instead of login (fix)
    return render_template('dashboard.html', user=session['user'], role=session['role'])

# ----------------------------------------------------------------------------------------
# Document Processing Function
def process_document(file):
    """Extracts and splits text from different document types."""
    try:
        ext = file.filename.split('.')[-1].lower()
        temp_path = tempfile.NamedTemporaryFile(delete=False).name  # Temporary file
        file.save(temp_path)

        documents = []

        #if ext == "txt":
            #loader = TextLoader(temp_path)
            #documents = loader.load()

        if ext == "pdf":
            loader = PyPDFLoader(temp_path)
            documents = loader.load()

        elif ext == "docx":
            doc = DocxDocument(temp_path)
            text = "\n".join([para.text for para in doc.paragraphs])
            documents = [Document(page_content=text)]

        elif ext == "xlsx":
            df = pd.read_excel(temp_path)
            text = "\n".join(df.astype(str).apply(lambda x: " | ".join(x), axis=1))
            documents = [Document(page_content=text)]

        elif ext == "pptx":
            prs = Presentation(temp_path)
            text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
            documents = [Document(page_content=text)]

        else:
            os.remove(temp_path)
            raise ValueError("Unsupported file format")

        os.remove(temp_path)  # Clean up temp file

        # Split text into smaller chunks
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
        split_docs = text_splitter.split_documents(documents)

        return split_docs

    except Exception as e:
        logger.error(f"Error processing document: {e}")
        raise

# ----------------------------------------------------------------------------------------
# File Upload (Restricted to Admins)
@app.route('/upload', methods=['POST'])
def upload_document():
    if 'user' not in session or session['role'] != 'admin':
        return jsonify({'error': 'Unauthorized access'}), 403

    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file provided'}), 400

        file = request.files['file']
        if file.filename == '' or not file.filename.endswith(('.pdf', '.docx', '.xlsx', '.pptx')):
            return jsonify({'error': 'Invalid file type'}), 400

        text_chunks = process_document(file)  # Extract and split text
        storage_service.upload_file(file, file.filename)  # Upload to S3
        vector_store.add_documents(text_chunks)  # Store in vector DB

        return jsonify({'message': 'File uploaded successfully'})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# ----------------------------------------------------------------------------------------
# User Query (Accessible to All Users)
@app.route('/query', methods=['POST'])
def query():
    if 'user' not in session:
        return jsonify({'error': 'Unauthorized access'}), 403

    data = request.json
    if 'question' not in data:
        return jsonify({'error': 'No question provided'}), 400

    try:
        response = llm_service.get_response(data['question'])
        return jsonify({'response': response})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# ----------------------------------------------------------------------------------------
# Delete Document (Restricted to Admins)
@app.route('/delete', methods=['POST'])
def delete_document():
    if 'user' not in session or session['role'] != 'admin':
        return jsonify({'error': 'Unauthorized access'}), 403

    data = request.json
    filename = data.get('filename')

    if not filename:
        return jsonify({'error': 'No filename provided'}), 400

    try:
        storage_service.delete_file(filename)
        vector_store.delete_documents(filename)  # delete vectors associated with the file.
        return jsonify({'message': f'File {filename} deleted successfully'})
    except Exception as e:
        return jsonify({'error': str(e)}), 500
##
# ----------------------------------------------------------------------------------------
# Run the web app
if __name__ == '__main__':
    app.run(host='0.0.0.0', port=8080, debug=True)
